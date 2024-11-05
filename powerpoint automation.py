import tkinter as tk
from tkinter import ttk
import requests
from pptx import Presentation
from pptx.util import Inches
import os

# Replace this with your actual Hugging Face API token
HUGGING_FACE_API_KEY = "YOUR_HUGGING_FACE_API_KEY"

# Define available languages
languages = [
    "English",
    "Bulgarian",
    "Spanish",
    "French",
    "German",
    "Italian",
    "Portuguese",
    "Russian",
    "Chinese",
    "Japanese"
]


def generate_slide_content(subject, slides, language="English"):
    slide_texts = []
    for i in range(1, slides + 1):
        prompt = f"Generate content in {language} for slide {i} for a presentation about '{subject} with information to each slide and point'."

        url = "https://api-inference.huggingface.co/models/EleutherAI/gpt-neo-2.7B"  # Replace with your chosen model's URL
        headers = {"Authorization": f"Bearer {HUGGING_FACE_API_KEY}"}
        data = {"inputs": prompt}

        response = requests.post(url, headers=headers, json=data)

        if response.status_code == 200:
            text = response.json()[0]['generated_text'].strip()
            slide_texts.append(text)
        else:
            print(f"Error {response.status_code}: {response.text}")
            slide_texts.append(f"Content could not be generated for slide {i}.")

    return slide_texts


def generate_slide_images(slides, subject, language="English"):
    images = []
    for i in range(1, slides + 1):
        prompt = f"Create an image for slide {i} of a presentation about '{subject}' in {language}."

        url = "https://api-inference.huggingface.co/models/CompVis/stable-diffusion-v1-4"  # Replace with your chosen image model's URL
        headers = {"Authorization": f"Bearer {HUGGING_FACE_API_KEY}"}
        data = {"inputs": prompt}

        response = requests.post(url, headers=headers, json=data)

        if response.status_code == 200:
            images.append(response.content)  # Stores the raw image data
        else:
            print(f"Error {response.status_code}: {response.text}")
            images.append(None)

    return images


def create_presentation(subject, slides, pictures, language):
    slide_texts = generate_slide_content(subject, slides, language)
    slide_images = generate_slide_images(slides, subject, language)

    # Create PowerPoint file
    pptx_folder = os.path.join(os.path.expanduser("~/Desktop/powerpoint"))
    os.makedirs(pptx_folder, exist_ok=True)  # Create the folder if it doesn't exist
    pptx_path = os.path.join(pptx_folder, f"{subject}.pptx")
    presentation = Presentation()

    for i in range(slides):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
        title = slide.shapes.title
        content = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        text_frame = content.text_frame

        # Set slide title and content
        title.text = f"Slide {i + 1}: {subject}"
        text_frame.add_paragraph().text = slide_texts[i]

        # (Optional) Add images to slide
        if slide_images[i]:
            # Save the image to a temporary file before inserting
            image_path = f"temp_image_{i}.png"
            with open(image_path, "wb") as img_file:
                img_file.write(slide_images[i])

            # Insert image into slide
            slide.shapes.add_picture(image_path, Inches(1), Inches(2), width=Inches(8))  # Adjust dimensions as needed

            # Clean up the temporary image file
            os.remove(image_path)

    # Save the PowerPoint presentation
    presentation.save(pptx_path)
    print(f"Presentation saved at {pptx_path}")


def create_gui():
    root = tk.Tk()
    root.title("PowerPoint Automation")

    # Label
    label = tk.Label(root, text="What would you like for your presentation?")
    label.pack(pady=10)

    # Text Entry for Subject
    subject_label = tk.Label(root, text="Subject:")
    subject_label.pack()
    subject_entry = tk.Entry(root)
    subject_entry.pack(pady=5)

    # Text Entry for Number of Slides
    slides_label = tk.Label(root, text="Number of Slides:")
    slides_label.pack()
    slides_entry = tk.Entry(root)
    slides_entry.pack(pady=5)

    # Text Entry for Number of Pictures
    pictures_label = tk.Label(root, text="Number of Pictures per Slide:")
    pictures_label.pack()
    pictures_entry = tk.Entry(root)
    pictures_entry.pack(pady=5)

    # Dropdown for Language Selection
    language_label = tk.Label(root, text="Select Language:")
    language_label.pack()
    selected_language = tk.StringVar()
    selected_language.set(languages[0])  # Default value
    language_dropdown = ttk.Combobox(root, textvariable=selected_language, values=languages)
    language_dropdown.pack(pady=5)

    # Submit Button
    submit_button = tk.Button(root, text="Generate Presentation", command=lambda: create_presentation(
        subject_entry.get(),
        int(slides_entry.get()),
        int(pictures_entry.get()),
        selected_language.get()
    ))
    submit_button.pack(pady=20)

    root.mainloop()


if __name__ == "__main__":
    create_gui()
